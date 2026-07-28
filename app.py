import streamlit as st
import pypdf
from docx import Document
from pptx import Presentation
import requests
import random
import json
import re
import base64

APP_NAME = "DocDigest"

st.set_page_config(page_title=APP_NAME, layout="centered", initial_sidebar_state="expanded")

# --- CUSTOM CSS FOR PERFECT ALIGNMENT & LUCIDE DESIGN LAYOUT ---
st.markdown("""
    <style>
        .flex-container {
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 10px;
        }
        .icon-svg {
            color: #6D28D9; /* Clean purple profile accent */
            flex-shrink: 0;
            vertical-align: middle;
        }
        .btn-flex {
            display: inline-flex;
            align-items: center;
            gap: 8px;
        }
        .memo-card {
            background-color: rgba(109, 40, 217, 0.05);
            padding: 16px 18px;
            border-radius: 8px;
            border: 1px solid rgba(109, 40, 217, 0.15);
            height: 100%;
        }
        section[data-testid="stSidebar"] .stButton > button {
            width: 100%;
            text-align: left;
            border-radius: 6px;
        }
        /* Lock the sidebar open — hide the collapse arrow entirely */
        [data-testid="stSidebarCollapseButton"] {
            display: none !important;
        }
        [data-testid="collapsedControl"] {
            display: none !important;
        }
    </style>
""", unsafe_allow_html=True)

# --- STATE MEMORY CORES: Keeps sections from wiping out ---
if "generated_summary" not in st.session_state:
    st.session_state.generated_summary = None
if "generated_notes" not in st.session_state:
    st.session_state.generated_notes = None
if "generated_sections" not in st.session_state:
    st.session_state.generated_sections = None
if "generated_mode_label" not in st.session_state:
    st.session_state.generated_mode_label = None
if "generated_batch_label" not in st.session_state:
    st.session_state.generated_batch_label = None
if "generated_cbt" not in st.session_state:
    st.session_state.generated_cbt = None
if "current_cbt_batch" not in st.session_state:
    st.session_state.current_cbt_batch = None
if "generated_cheatsheet" not in st.session_state:
    st.session_state.generated_cheatsheet = None
if "active_view" not in st.session_state:
    st.session_state.active_view = None  # "analogy" | "quiz" | "cheatsheet"
if "transcribed_notes" not in st.session_state:
    st.session_state.transcribed_notes = None

# --- TITLE WITH PURE LUCIDE BOOKSTACK ---
st.markdown(f"""
    <div class="flex-container">
        <svg class="icon-svg" xmlns="http://www.w3.org/2000/svg" width="36" height="36" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round">
            <path d="m16 6 4 14"/>
            <path d="M12 6v14"/>
            <path d="M8 8v12"/>
            <path d="M4 4v16"/>
        </svg>
        <h1 style="margin: 0; padding: 0; font-size: 2.2rem;">{APP_NAME}</h1>
    </div>
""", unsafe_allow_html=True)

st.write("Upload your lecture notes, slides, or PDFs, then choose how you want to conquer them.")


def extract_text(uploaded_file):
    filename = uploaded_file.name
    text = ""
    try:
        if filename.endswith('.pdf'):
            pdf_reader = pypdf.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                try:
                    t = page.extract_text()
                    if t: text += t + "\n"
                except: continue
        elif filename.endswith('.docx'):
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif filename.endswith('.pptx') or filename.endswith('.pptm'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"
    except Exception as e:
        return f"Error reading file structure: {e}"
    return text


def ask_gemini(api_key, prompt_text, dynamic_mode=False):
    models = ["gemini-2.5-flash", "gemini-1.5-flash"]
    generation_config = {"temperature": 0.85 if dynamic_mode else 0.2}

    for model in models:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
        headers = {'Content-Type': 'application/json'}
        payload = {
            "contents": [{"parts": [{"text": prompt_text}]}],
            "generationConfig": generation_config
        }
        try:
            response = requests.post(url, headers=headers, json=payload)
            response_json = response.json()
            if 'candidates' in response_json and response_json['candidates']:
                return response_json['candidates'][0]['content']['parts'][0]['text']
            elif 'error' in response_json:
                error_msg = response_json['error']['message']
                if "demand" in error_msg.lower() or "quota" in error_msg.lower() or "not found" in error_msg.lower():
                    continue
                return f"Google API Error: {error_msg}"
        except: continue

    # --- PURE LUCIDE ALERT BOX ---
    return """
    <div style="display:flex; align-items:center; gap:8px; color:#DC2626; font-weight:600;">
        <svg xmlns="http://www.w3.org/2000/svg" width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="m21.73 18-8-14a2 2 0 0 0-3.48 0l-8 14A2 2 0 0 0 4 21h16a2 2 0 0 0 1.73-3Z"/><line x1="12" y1="9" x2="12" y2="13"/><line x1="12" y1="17" x2="12.01" y2="17"/></svg>
        Server lines are busy. Tap the button again!
    </div>
    """


def transcribe_images(api_key, image_files):
    """Sends photos directly to Gemini's multimodal endpoint and asks it to
    transcribe the handwritten/printed content into plain text, in the order
    given. This deliberately skips OCR libraries like Tesseract — Gemini's
    vision understanding handles messy handwriting far better."""
    models = ["gemini-2.5-flash", "gemini-1.5-flash"]

    parts = []
    for img_file in image_files:
        img_bytes = img_file.getvalue()
        mime_type = img_file.type or "image/jpeg"
        encoded = base64.b64encode(img_bytes).decode("utf-8")
        parts.append({"inline_data": {"mime_type": mime_type, "data": encoded}})

    parts.append({
        "text": (
            "Transcribe the handwritten and/or printed content of these images "
            "verbatim into plain text, in the exact order the images were given "
            "(treat them as consecutive pages). Preserve structure where visible "
            "(headings, bullet points, numbered lists). Do not summarize, "
            "correct, or add commentary — output ONLY the transcribed text."
        )
    })

    for model in models:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
        headers = {'Content-Type': 'application/json'}
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": {"temperature": 0.1}
        }
        try:
            response = requests.post(url, headers=headers, json=payload)
            response_json = response.json()
            if 'candidates' in response_json and response_json['candidates']:
                return response_json['candidates'][0]['content']['parts'][0]['text']
            elif 'error' in response_json:
                error_msg = response_json['error']['message']
                if "demand" in error_msg.lower() or "quota" in error_msg.lower() or "not found" in error_msg.lower():
                    continue
                return None
        except Exception:
            continue
    return None


def extract_json_block(text):
    """Only used for the paired notes/analogy feature, since that one needs
    structured per-item output to render as aligned rows."""
    if not text:
        return None
    cleaned = re.sub(r"```json|```", "", text).strip()
    match = re.search(r"(\[.*\]|\{.*\})", cleaned, re.DOTALL)
    if not match:
        return None
    try:
        return json.loads(match.group(1))
    except Exception:
        return None


# --- SMART API KEY ROTATION SETUP (kept in the sidebar, same as original) ---
api_keys = []
if "GEMINI_API_KEY_1" in st.secrets and st.secrets["GEMINI_API_KEY_1"]:
    api_keys.append(st.secrets["GEMINI_API_KEY_1"])
if "GEMINI_API_KEY_2" in st.secrets and st.secrets["GEMINI_API_KEY_2"]:
    api_keys.append(st.secrets["GEMINI_API_KEY_2"])
if "GEMINI_API_KEY_3" in st.secrets and st.secrets["GEMINI_API_KEY_3"]:
    api_keys.append(st.secrets["GEMINI_API_KEY_3"])

manual_key = st.sidebar.text_input("Backup API Key Entry (Optional)", type="password")
if manual_key:
    api_keys.append(manual_key)

api_key = random.choice(api_keys) if api_keys else None

# --- FILE UPLOAD STAYS IN THE MAIN AREA ---
uploaded_file = st.file_uploader("Drop your study document here (PDF, DOCX, PPTX, PPTM)", type=["pdf", "docx", "pptx", "pptm"])

st.markdown("**— or —**")

uploaded_images = st.file_uploader(
    "Upload photos of handwritten notes (multiple pages, in page order)",
    type=["png", "jpg", "jpeg"],
    accept_multiple_files=True
)
if uploaded_images:
    st.caption(f"{len(uploaded_images)} photo(s) selected. Larger batches take a bit longer and use more of your API quota.")
    if st.button("Transcribe These Photos"):
        if not api_key:
            st.error("Missing API Key!")
        else:
            with st.spinner("Reading your handwritten pages..."):
                result = transcribe_images(api_key, uploaded_images)
                if result:
                    st.session_state.transcribed_notes = result
                else:
                    st.error("Couldn't transcribe these photos — please try again.")

has_content = bool(uploaded_file or st.session_state.transcribed_notes)

if has_content:
    if uploaded_file:
        raw_text = extract_text(uploaded_file)
        source_label = uploaded_file.name
    else:
        raw_text = st.session_state.transcribed_notes
        source_label = "Transcribed photos"

    st.caption(f"📄 Active source: {source_label}")

    total_length = len(raw_text) if raw_text else 0
    chunk_size = max(1, total_length // 4)

    chunk_1 = raw_text[0:chunk_size] if total_length > 0 else ""
    chunk_2 = raw_text[chunk_size:chunk_size*2] if total_length > 0 else ""
    chunk_3 = raw_text[chunk_size*2:chunk_size*3] if total_length > 0 else ""
    chunk_4 = raw_text[chunk_size*3:] if total_length > 0 else ""

    # ============================================================
    # SIDEBAR — every control clustered here, locked open
    # ============================================================
    with st.sidebar:
        st.markdown("""
            <div class="flex-container">
                <svg class="icon-svg" xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M2 3h6a4 4 0 0 1 4 4v14a3 3 0 0 0-3-3H2z"/><path d="M22 3h-6a4 4 0 0 0-4 4v14a3 3 0 0 1 3-3h7z"/></svg>
                <h3 style="margin:0;">Tailored Multi-Mode Explanation Engine</h3>
            </div>
        """, unsafe_allow_html=True)

        st.caption("Pick a persona — both columns follow the same numbered structure, so you can read them in parallel.")

        safe_combined_text = (
            f"[Introductory Section]\n{chunk_1[:3500]}\n\n"
            f"[Core Section A]\n{chunk_2[:3500]}\n\n"
            f"[Core Section B]\n{chunk_3[:3500]}\n\n"
            f"[Advanced / Concluding Section]\n{chunk_4[:3500]}"
        )

        MODE_OPTIONS = {
            "buddy": ("Campus Buddy", "You are a witty, supportive university peer tutor writing a cohesive campus survival guide. Establish a single, continuous story using the chaotic ecosystem of university life (like navigating strict hostel porters, fighting for hostel Wi-Fi bandwidth, or queuing at the campus gate). Every technical concept from the document must become a recurring character or mechanic in this campus saga. Keep the tone conversational, funny, and deeply protective of your friend's grades. Wrap key terms in <strong> tags."),
            "street": ("Street-Smart", "You are a practical operations mentor grounding the material in a single, continuous real-world hustle. Use a macro-metaphor of a massive physical enterprise—like a city-wide delivery logistics network, a bustling open-air market, or an operation tracking down counterfeit goods. Ensure every piece of technical infrastructure on the left is represented as a structural asset or security check in this ongoing street-smart operation. Keep it highly interconnected, punchy, and grounded. Wrap key terms in <strong> tags."),
            "technical": ("Deep Technical", "You are a senior enterprise systems architect designing a fully integrated dependency pipeline. Do not just list isolated specs; write a continuous architectural whitepaper where every technical concept acts as a precise mechanism that feeds into, triggers, or constrains the next phase of the global infrastructure layout. Ensure deep logical continuity from top to bottom. Wrap key terms in <strong> tags."),
            "layman": ("Layman", "You are a master storyteller translating complex machinery into a beautiful, non-technical visual fable. Use a single, expansive overarching metaphor—like a giant castle postal sorting system, a kingdom's plumbing infrastructure, or a massive restaurant kitchen. Walk the reader through this world continuously, ensuring that each new concept builds directly upon the room or process established in the previous paragraph. Wrap key terms in <strong> tags."),
            "naija": ("Naija Voice", """You are the "Naija Voice" persona: convert the technical content into a relatable Nigerian street/market/society analogy — using scenarios like market bargaining, "one chance" buses, NEPA/generator light issues, okada riders, danfo conductors, network wahala, POS agents, and everyday Naija social dynamics (landlord issues, "aunty" at the shop, etc).

Internally, before writing (do not show this reasoning in your output): (1) extract every key technical term with its precise one-line definition, (2) map each term to an exact Naija-world equivalent that preserves the SAME underlying mechanism, not just a similar vibe or mood — if a term has no natural match, don't force one, (3) narrate using ONLY those confirmed mappings, (4) self-check for anywhere the analogy could mislead a reader about how the real concept actually works.

Write the visible narrative in a warm, expressive Nigerian pidgin-inflected voice (natural code-switching between English and pidgin — don't overdo the pidgin to the point of losing clarity). If your internal self-check found a real gap, weave a short honest caveat naturally into the end of the narrative instead of dropping it silently.

Tone rules: confident, expressive, funny where natural — but never mock or stereotype; write like someone proudly explaining to a younger sibling or junior colleague. Avoid inventing slang that doesn't map to a real concept just to sound "street." Don't over-explain the pidgin — trust the reader.

Wrap key terms in <strong> tags."""),
        }

        for mode_id, (label, persona_vibe) in MODE_OPTIONS.items():
            if st.button(label, key=f"mode_{mode_id}"):
                if not api_key:
                    st.error("Missing API Key!")
                elif not raw_text.strip():
                    st.error("Could not extract any text from this document.")
                else:
                    with st.spinner("Building the dual-stream notes..."):
                        prompt = f"""
                        You are a dual-stream content engine. Generate two perfectly distinct but
                        structurally mirrored texts that will sit side-by-side.

                        FIRST, break the document down into a numbered sequence of concepts/sections
                        in the order they appear (e.g. 1, 2, 3...). Use as many numbered sections as the
                        document naturally has distinct concepts — do not force an arbitrary count.

                        For EACH section, produce a matched pair:
                        - "literal_note": a highly accurate, structured, literal academic explanation
                          of THIS section only — definitions, configurations, core technical rules.
                        - "persona_note": the SAME section re-explained through the chosen Persona below.
                          Do NOT write an isolated, dry, one-off translation — write it as a piece of
                          one continuous "Companion Textbook", referencing and building on earlier
                          sections' analogies where natural, so reading all persona_note fields in
                          order feels like one flowing story, not disconnected snippets.

                        PERSONA INSTRUCTION FOR "persona_note":
                        {persona_vibe}

                        STIPULATION ON BOLDING: Never use markdown asterisks (like **text**) to bold
                        phrases. Instead, wrap key technical definitions and core takeaways using HTML
                        strong tags (like <strong>text</strong>).

                        Respond with ONLY a raw JSON array, no markdown fences, no commentary, in this
                        exact schema:
                        [
                          {{
                            "number": 1,
                            "title": "short section title, e.g. 'Sender'",
                            "literal_note": "...",
                            "persona_note": "..."
                          }}
                        ]

                        Document Text:
                        {safe_combined_text}
                        """
                        result = ask_gemini(api_key, prompt, dynamic_mode=True)
                        parsed = extract_json_block(result)
                        if parsed:
                            st.session_state.generated_sections = parsed
                            # Rebuild flat text blobs too, so the Quiz and Recall Hook Table
                            # (which expect one continuous passage) keep working unchanged.
                            st.session_state.generated_notes = "\n\n".join(
                                f"{s.get('number', '')}. {s.get('title', '')}\n{s.get('literal_note', '')}"
                                for s in parsed
                            )
                            st.session_state.generated_summary = "\n\n".join(
                                f"{s.get('number', '')}. {s.get('title', '')}\n{s.get('persona_note', '')}"
                                for s in parsed
                            )
                            st.session_state.generated_mode_label = label
                            st.session_state.active_view = "analogy"
                        else:
                            st.error("Couldn't generate the dual-stream notes this time — please try again.")
                            with st.expander("See raw response (for debugging)"):
                                st.code(result or "(no response — API call itself failed)")

        st.markdown("---")
        st.markdown("""
            <div class="flex-container">
                <svg class="icon-svg" xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M12 5a3 3 0 1 0-5.997.125 4 4 0 0 0-2.526 5.77 4 4 0 0 0 .556 6.588A4 4 0 1 0 12 18Z"/><path d="M12 5a3 3 0 1 1 5.997.125 4 4 0 0 1 2.526 5.77 4 4 0 0 1-.556 6.588A4 4 0 1 1 12 18Z"/><path d="M12 5v14"/><path d="M12 12h6"/><path d="M12 12H6"/><path d="M12 7h5"/><path d="M12 16h5"/><path d="M12 7H7"/><path d="M12 16H7"/></svg>
                <h3 style="margin:0;">Theory-to-CBT Objective Drill</h3>
            </div>
        """, unsafe_allow_html=True)

        st.caption("Tests both the real technical facts AND whether you understood the analogy — built from the analogy currently shown above.")

        has_analogy = bool(st.session_state.generated_notes and st.session_state.generated_summary)

        if not has_analogy:
            st.info("Generate a persona analogy above first — the quiz is built from it.")

        quiz_batch = st.selectbox(
            "Focus the quiz on:",
            ["Batch 1: Introduction & Foundation Concepts",
             "Batch 2: Core Methodologies & Process Details",
             "Batch 3: Deep Technical Content",
             "Batch 4: Advanced Scenarios & Conclusions"],
            key="quiz_batch_select"
        )

        if quiz_batch.startswith("Batch 1"):
            quiz_focus_text = chunk_1[:9000]
        elif quiz_batch.startswith("Batch 2"):
            quiz_focus_text = chunk_2[:9000]
        elif quiz_batch.startswith("Batch 3"):
            quiz_focus_text = chunk_3[:9000]
        else:
            quiz_focus_text = chunk_4[:9000]

        if st.button("Generate Analogy-Aware Quiz", key="btn_cbt", disabled=not has_analogy):
            if not api_key:
                st.error("Missing API Key!")
            elif not quiz_focus_text.strip():
                st.error("This section doesn't contain enough text to extract data.")
            else:
                with st.spinner("Writing analogy-aware quiz questions..."):
                    prompt = f"""
                    You are generating CBT-style multiple choice questions that test BOTH technical
                    accuracy AND correct understanding of the persona analogy the user just read.

                    Only generate questions covering concepts that actually appear in the "Section To
                    Quiz" text below. Use the full Grounded Truth and Persona Narrative only as
                    reference material to find the correct definitions and analogy mappings for those
                    concepts — ignore concepts from other parts of the document that aren't present
                    in this section.

                    Step 1 — Identify the core technical concept(s) covered in the Section To Quiz.
                    Step 2 — Identify the exact analogy mapping used in the Persona Narrative for each concept.
                    Step 3 — For each question:
                    - Write ONE correct answer that reflects the accurate mapping from Step 2.
                    - Write ONE distractor that represents the most common way someone misreads THIS
                      SPECIFIC analogy (a real, plausible misunderstanding — not a random wrong fact).
                    - Write 1-2 more distractors that are plausible but clearly wrong on technical grounds.
                    - After the answer, add a one-line explanation of WHY the misread distractor is a
                      common trap, so the user learns even from getting it wrong.
                    Step 4 — Self-check internally: confirm no question's "correct" answer actually
                    reflects a subtly wrong mapping (re-verify against the Grounded Truth directly, not
                    the analogy). Do not show this step in your output.

                    Generate 6 to 8 questions total.

                    Output ONLY this format per question, no internal steps shown:
                    Question: [text]
                    A) [option] B) [option] C) [option] D) [option]
                    ✅ Correct Answer: [letter] — [1-line technical justification, quoting the Grounded Truth's logic, not the analogy]
                    ⚠️ Common trap: [which distractor people usually pick and why the analogy makes that mistake tempting]

                    Section To Quiz ({quiz_batch}):
                    {quiz_focus_text}

                    Full Grounded Truth (reference only):
                    {st.session_state.generated_notes}

                    Full Persona Narrative used ({st.session_state.generated_mode_label}, reference only):
                    {st.session_state.generated_summary}
                    """
                    st.session_state.generated_cbt = ask_gemini(api_key, prompt, dynamic_mode=False)
                    st.session_state.current_cbt_batch = f"{st.session_state.generated_mode_label} — {quiz_batch}"
                    st.session_state.active_view = "quiz"

        st.markdown("---")
        st.markdown("""
            <div class="flex-container">
                <svg class="icon-svg" xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M12 3v18"/><path d="M3 12h18"/><rect x="3" y="3" width="18" height="18" rx="2"/></svg>
                <h3 style="margin:0;">Recall Hook Table</h3>
            </div>
        """, unsafe_allow_html=True)

        st.caption("Anchors each term to the specific analogy moment that explained it — built from the analogy currently shown above.")

        if not has_analogy:
            st.info("Generate a persona analogy above first — the table is built from it.")

        if st.button("Generate Recall Hook Table", key="btn_table", disabled=not has_analogy):
            if not api_key:
                st.error("Missing API Key!")
            else:
                with st.spinner("Anchoring each term to its analogy hook..."):
                    prompt = f"""
                    You are generating a compact recall table that anchors each technical term to
                    the analogy hook that helped the user understand it — not just a flat definition.

                    Step 1 — List every key term from the Deep Technical passage with its precise definition.
                    Step 2 — For each term, pull the exact phrase or moment from the persona narrative
                    that represents that term (quote or closely paraphrase the specific analogy line —
                    not a generic restatement).
                    Step 3 — Self-check internally: confirm the recall hook for each term matches Step
                    1's definition exactly — if the persona narrative's phrasing has drifted from the
                    technical meaning, fix the hook to be accurate rather than copying the drifted
                    version verbatim. Do not show this step in your output.

                    Output ONLY a markdown table with these columns (no internal steps shown):
                    | Term | Definition | Recall Hook ({st.session_state.generated_mode_label}) |

                    Keep each Recall Hook under 20 words — punchy, memorable, not a full re-explanation.

                    Deep Technical passage:
                    {st.session_state.generated_notes}

                    Persona narrative used:
                    {st.session_state.generated_summary}
                    """
                    st.session_state.generated_cheatsheet = ask_gemini(api_key, prompt, dynamic_mode=False)
                    st.session_state.active_view = "cheatsheet"

    # ============================================================
    # MAIN AREA — presentation of whichever result is active
    # ============================================================
    view = st.session_state.active_view

    if view is None:
        st.info("Document loaded. Choose an action from the left panel to get started.")

    elif view == "analogy":
        st.markdown(f"### 💡 {st.session_state.generated_mode_label} — Dual-Stream Notes")
        sections = st.session_state.generated_sections or []

        for sec in sections:
            number = sec.get("number", "")
            title = sec.get("title", "").strip()
            note = sec.get("literal_note", "").strip()
            persona_note = sec.get("persona_note", "").strip()

            st.markdown(f"#### {number}. {title}")
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**📌 Grounded Truth**")
                st.markdown(f'<div class="memo-card">{note or "—"}</div>', unsafe_allow_html=True)
            with c2:
                st.markdown("**💡 Immersive Persona Note**")
                st.markdown(f'<div class="memo-card">{persona_note or "—"}</div>', unsafe_allow_html=True)
            st.markdown("")  # spacer between sections

    elif view == "quiz":
        st.markdown(f"### 🧩 Analogy-Aware Quiz — {st.session_state.current_cbt_batch}")
        if st.session_state.generated_cbt:
            st.markdown(st.session_state.generated_cbt, unsafe_allow_html=True)

    elif view == "cheatsheet":
        st.markdown(f"### 🗂️ Recall Hook Table — {st.session_state.generated_mode_label}")
        if st.session_state.generated_cheatsheet:
            st.markdown(st.session_state.generated_cheatsheet, unsafe_allow_html=True)

else:
    st.info(f"Upload a document above to unlock {APP_NAME}'s tools.")
